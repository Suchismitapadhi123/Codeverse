import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Color Scheme Constants
BG_COLOR = RGBColor(12, 10, 22)       # Dark Obsidian
TEAL_COLOR = RGBColor(0, 242, 254)    # Neon Teal (Accent 1)
GREEN_COLOR = RGBColor(0, 255, 135)   # Cyber Green (Accent 2)
VIOLET_COLOR = RGBColor(157, 78, 237) # Neon Violet (Accent 3)
WHITE_COLOR = RGBColor(255, 255, 255) # Text Primary
GRAY_COLOR = RGBColor(180, 180, 190)  # Text Secondary

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_slide_title(slide, text, color=TEAL_COLOR):
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Consolas"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = color
    return txBox

# ----------------- SLIDE 1: Title Slide -----------------
slide_layout = prs.slide_layouts[6] # Blank
slide1 = prs.slides.add_slide(slide_layout)
set_slide_background(slide1)

# Large Title
title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.0))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "CodeVerse"
p.font.name = "Consolas"
p.font.size = Pt(72)
p.font.bold = True
p.font.color.rgb = TEAL_COLOR
p.alignment = PP_ALIGN.CENTER

# Glow separator bar
# (python-pptx doesn't support easy lines, we use dashes)
p2 = tf.add_paragraph()
p2.text = "=================================================="
p2.font.name = "Consolas"
p2.font.size = Pt(18)
p2.font.color.rgb = GREEN_COLOR
p2.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide1.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(11.3), Inches(1.5))
tf_sub = subtitle_box.text_frame
tf_sub.word_wrap = True
p_sub = tf_sub.paragraphs[0]
p_sub.text = "A Premium Cyberpunk Student Programming Portal & Learning Console"
p_sub.font.name = "Consolas"
p_sub.font.size = Pt(20)
p_sub.font.bold = True
p_sub.font.color.rgb = WHITE_COLOR
p_sub.alignment = PP_ALIGN.CENTER

p_sub2 = tf_sub.add_paragraph()
p_sub2.text = "Designed for College Students • Dynamic Tracks • MCQ Quizzes • 3D Flashcards • Free Live AI"
p_sub2.font.name = "Consolas"
p_sub2.font.size = Pt(14)
p_sub2.font.color.rgb = GRAY_COLOR
p_sub2.alignment = PP_ALIGN.CENTER


# ----------------- SLIDE 2: Project Overview -----------------
slide2 = prs.slides.add_slide(slide_layout)
set_slide_background(slide2)
add_slide_title(slide2, "> Project Overview")

content_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf2 = content_box.text_frame
tf2.word_wrap = True

bullets2 = [
    ("Target Audience", "Specially designed for college programming and engineering students.", TEAL_COLOR),
    ("Theme Design", "Glowing monospaced typography, obsidian terminal layouts, and CRT-glowing borders for high visual engagement.", GREEN_COLOR),
    ("Gamification", "Integrates daily active commit streaks and Circular SVG completion progress dials to drive student retention.", VIOLET_COLOR),
    ("Core Methodology", "Builds active recall memory and syntax familiarity through MCQ Output Prediction Quizzes and 3D concept card decks.", WHITE_COLOR)
]

for title, desc, col in bullets2:
    p = tf2.add_paragraph()
    p.space_after = Pt(14)
    run_title = p.add_run()
    run_title.text = f"■  {title}: "
    run_title.font.name = "Consolas"
    run_title.font.bold = True
    run_title.font.size = Pt(20)
    run_title.font.color.rgb = col
    
    run_desc = p.add_run()
    run_desc.text = desc
    run_desc.font.name = "Consolas"
    run_desc.font.size = Pt(18)
    run_desc.font.color.rgb = WHITE_COLOR


# ----------------- SLIDE 3: System Architecture -----------------
slide3 = prs.slides.add_slide(slide_layout)
set_slide_background(slide3)
add_slide_title(slide3, "> Tech Stack & Architecture")

content_box3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf3 = content_box3.text_frame
tf3.word_wrap = True

bullets3 = [
    ("Frontend Layer", "Structured using semantic HTML5, stylized using Vanilla CSS (glowing cyberpunk border cards, glassmorphism, responsive navigation), and wired using modular ES6 JavaScript."),
    ("State Management", "Maintains developer streaks, quiz grade arrays, badge achievements, and mastered cards inside the browser's localStorage for instant offline access and persistence."),
    ("Local Web Server", "Runs locally on port 8081 via Python's built-in HTTP server wrapper, supporting relative PDF links and native browser rendering."),
    ("Data Layer", "Three dynamic JS generators parse core templates and expand data tables to exactly 100 entries per language track at runtime, preventing file bloat.")
]

for title, desc in bullets3:
    p = tf3.add_paragraph()
    p.space_after = Pt(16)
    r1 = p.add_run()
    r1.text = f"■  {title}\n"
    r1.font.name = "Consolas"
    r1.font.bold = True
    r1.font.size = Pt(18)
    r1.font.color.rgb = GREEN_COLOR
    
    r2 = p.add_run()
    r2.text = f"   {desc}"
    r2.font.name = "Consolas"
    r2.font.size = Pt(16)
    r2.font.color.rgb = GRAY_COLOR


# ----------------- SLIDE 4: Interactive Tracks & Notes -----------------
slide4 = prs.slides.add_slide(slide_layout)
set_slide_background(slide4)
add_slide_title(slide4, "> Study Tracks & Notes Integration")

content_box4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf4 = content_box4.text_frame
tf4.word_wrap = True

bullets4 = [
    ("Dynamic Language Tracks", "Covers four major college coding tracks: Python, JavaScript, C++, and Algorithms & DSA."),
    ("Dynamic UI Themes", "The sidebar and panels dynamically switch styling accents to coordinate with track branding (Teal for Python, Amber for JavaScript, Violet for C++, Green for DSA)."),
    ("Syllabus Notes Toggle", "A premium dual-tab selector enables students to toggle between:"),
    ("   - Tab A: Interactive Docs", "Loads summarized outlines from Wikipedia API and recommended textbooks from Open Library API."),
    ("   - Tab B: Original Handbooks", "Embeds the professor's PDF note files ('The Ultimate Python Handbook.pdf', 'JS_Chapterwise_Notes.pdf', etc.) directly inside an iframe terminal window, with single-click download actions.")
]

for title, desc in bullets4:
    p = tf4.add_paragraph()
    p.space_after = Pt(10)
    r1 = p.add_run()
    r1.text = f"■  {title}: " if not title.startswith("   -") else f"      {title}: "
    r1.font.name = "Consolas"
    r1.font.bold = True
    r1.font.size = Pt(18)
    r1.font.color.rgb = VIOLET_COLOR if not title.startswith("   -") else GREEN_COLOR
    
    r2 = p.add_run()
    r2.text = desc
    r2.font.name = "Consolas"
    r2.font.size = Pt(16)
    r2.font.color.rgb = WHITE_COLOR if not title.startswith("   -") else GRAY_COLOR


# ----------------- SLIDE 5: 3D Flashcard Console -----------------
slide5 = prs.slides.add_slide(slide_layout)
set_slide_background(slide5)
add_slide_title(slide5, "> 3D Active Recall Flashcards")

content_box5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf5 = content_box5.text_frame
tf5.word_wrap = True

bullets5 = [
    ("3D CSS Perspectives", "Holographic cards flip 180 degrees using CSS 3D matrix properties (backface-visibility hidden) for a premium tactile feel."),
    ("Keyboard Hotkeys", "Supports active keyboard event listener: Spacebar to flip card, Left/Right Arrows to navigate back and forth."),
    ("Mastered Tracker", "Allows students to label cards as 'Mastered' which saves to their profile, logs streaks, and fires visual confetti bursts."),
    ("Scaled Database", "Expanded to exactly 100 concepts per track (400 total cards) using unique incremental IDs and customized compiler validation tips on the card backs.")
]

for title, desc in bullets5:
    p = tf5.add_paragraph()
    p.space_after = Pt(16)
    r1 = p.add_run()
    r1.text = f"■  {title}\n"
    r1.font.name = "Consolas"
    r1.font.bold = True
    r1.font.size = Pt(18)
    r1.font.color.rgb = TEAL_COLOR
    
    r2 = p.add_run()
    r2.text = f"   {desc}"
    r2.font.name = "Consolas"
    r2.font.size = Pt(16)
    r2.font.color.rgb = GRAY_COLOR


# ----------------- SLIDE 6: Live Free AI Doubt Solver -----------------
slide6 = prs.slides.add_slide(slide_layout)
set_slide_background(slide6)
add_slide_title(slide6, "> Live Free AI Doubt Solver")

content_box6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf6 = content_box6.text_frame
tf6.word_wrap = True

bullets6 = [
    ("Keyless Integration", "Wired with Pollinations.ai's free public inference text API. No API keys or registration required."),
    ("CORS-Enabled", "Allows direct browser-to-API requests completely on the client-side, avoiding CORS blockages that occur on Anthropic or OpenAI endpoints."),
    ("Interactive Sidebar", "Holds 100 pre-saved exam-asked questions filterable by keyword search and categorized language tags (Python, JS, C++, DSA)."),
    ("Markdown Parsing Console", "Converts raw text streams into rich HTML, displaying code blocks, bold text, lists, and dynamic markdown blockquote notes for student tips.")
]

for title, desc in bullets6:
    p = tf6.add_paragraph()
    p.space_after = Pt(16)
    r1 = p.add_run()
    r1.text = f"■  {title}\n"
    r1.font.name = "Consolas"
    r1.font.bold = True
    r1.font.size = Pt(18)
    r1.font.color.rgb = GREEN_COLOR
    
    r2 = p.add_run()
    r2.text = f"   {desc}"
    r2.font.name = "Consolas"
    r2.font.size = Pt(16)
    r2.font.color.rgb = GRAY_COLOR


# ----------------- SLIDE 7: Analytics & Streaks -----------------
slide7 = prs.slides.add_slide(slide_layout)
set_slide_background(slide7)
add_slide_title(slide7, "> Developer Metrics & Streaks")

content_box7 = slide7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf7 = content_box7.text_frame
tf7.word_wrap = True

bullets7 = [
    ("Responsive Stats Row", "Top-level summary dashboard with 3 cards: Daily Active Streak, Total Topics Compiled, and Flashcards Mastered."),
    ("Circular SVG Dials", "Renders beautiful circular gauges tracing percentage completions of syllabus modules per language track."),
    ("Practice Diagnostics", "Vertical column bar charts illustrating student average scores obtained in output prediction quizzes."),
    ("Compiler Milestones", "Locked/Unlocked achievement badges (First Compile, Systems Engineer, Card Scholar, Card Master) that unlock and trigger toast notifications.")
]

for title, desc in bullets7:
    p = tf7.add_paragraph()
    p.space_after = Pt(16)
    r1 = p.add_run()
    r1.text = f"■  {title}\n"
    r1.font.name = "Consolas"
    r1.font.bold = True
    r1.font.size = Pt(18)
    r1.font.color.rgb = VIOLET_COLOR
    
    r2 = p.add_run()
    r2.text = f"   {desc}"
    r2.font.name = "Consolas"
    r2.font.size = Pt(16)
    r2.font.color.rgb = GRAY_COLOR


# ----------------- SLIDE 8: Summary / Conclusion -----------------
slide8 = prs.slides.add_slide(slide_layout)
set_slide_background(slide8)
add_slide_title(slide8, "> Conclusion", color=GREEN_COLOR)

conclusion_box = slide8.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(4.5))
tf8 = conclusion_box.text_frame
tf8.word_wrap = True

p = tf8.paragraphs[0]
p.text = "CodeVerse delivers a premium, cohesive, and feature-rich educational platform that bridges offline syllabus materials and live online AI capabilities."
p.font.name = "Consolas"
p.font.size = Pt(22)
p.font.color.rgb = WHITE_COLOR
p.space_after = Pt(24)

p2 = tf8.add_paragraph()
p2.text = "✓ 400 MCQ Questions Database\n✓ 400 Concept Flashcards Deck\n✓ 400 Doubts & Interview Presets\n✓ Multi-Page State Persistence\n✓ 100% Free Live Online AI Tutor Integration"
p2.font.name = "Consolas"
p2.font.size = Pt(20)
p2.font.bold = True
p2.font.color.rgb = TEAL_COLOR
p2.space_after = Pt(24)

p3 = tf8.add_paragraph()
p3.text = "[SYSTEM OUTPUT]: Deployment Successful. Ready for local compilation."
p3.font.name = "Consolas"
p3.font.size = Pt(16)
p3.font.color.rgb = GREEN_COLOR

# Save the Presentation
output_path = os.path.join(r"C:\Users\SUCHISMITA\.gemini\antigravity\scratch\codeverse", "CodeVerse_Presentation.pptx")
prs.save(output_path)
print("Presentation saved successfully at:", output_path)
