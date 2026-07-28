import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# --- Color Scheme: Purple Dark & Green Dark Theme ---
BG_COLOR = RGBColor(10, 8, 20)          # Very Dark Obsidian Purple
CARD_BG_COLOR = RGBColor(22, 18, 38)     # Slightly Lighter Dark Purple
BORDER_COLOR = RGBColor(155, 93, 229)    # Neon Purple/Violet (Primary Accent)
GREEN_COLOR = RGBColor(0, 245, 150)      # Cyber Green (Secondary Accent)
WHITE_COLOR = RGBColor(255, 255, 255)    # Text Primary
GRAY_COLOR = RGBColor(175, 172, 190)     # Text Secondary

# Helper function to set slide background
def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

# Helper function to draw header with a glowing line
def add_slide_header(slide, title_text):
    # Title Text Box
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Consolas"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = GREEN_COLOR
    
    # Glowing line under title (thin filled rectangle)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0.8), Inches(1.1), Inches(11.73), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER_COLOR
    line.line.fill.background() # No border line

# Helper function to draw a card panel
def draw_card(slide, left, top, width, height, num_str, title_str, body_str):
    # Base Card Shape
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG_COLOR
    card.line.color.rgb = BORDER_COLOR
    card.line.width = Pt(1.5)
    
    # Text Box overlay for formatting
    tx_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    # Header Paragraph: "01   Study Tracks"
    p_header = tf.paragraphs[0]
    p_header.space_after = Pt(12)
    
    # Number
    run_num = p_header.add_run()
    run_num.text = f"{num_str}   "
    run_num.font.name = "Consolas"
    run_num.font.bold = True
    run_num.font.size = Pt(16)
    run_num.font.color.rgb = BORDER_COLOR
    
    # Title
    run_title = p_header.add_run()
    run_title.text = title_str
    run_title.font.name = "Consolas"
    run_title.font.bold = True
    run_title.font.size = Pt(18)
    run_title.font.color.rgb = WHITE_COLOR
    
    # Divider line under title inside card
    p_div = tf.add_paragraph()
    p_div.space_after = Pt(12)
    run_div = p_div.add_run()
    run_div.text = "-" * 32
    run_div.font.name = "Consolas"
    run_div.font.size = Pt(10)
    run_div.font.color.rgb = RGBColor(60, 50, 80)
    
    # Body Paragraph
    p_body = tf.add_paragraph()
    run_body = p_body.add_run()
    run_body.text = body_str
    run_body.font.name = "Consolas"
    run_body.font.size = Pt(13)
    run_body.font.color.rgb = GRAY_COLOR

# ----------------- SLIDE 1: Title Slide -----------------
slide_layout = prs.slide_layouts[6] # Blank
slide1 = prs.slides.add_slide(slide_layout)
set_slide_background(slide1)

# Large Title
title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.0))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "CodeVerse"
p.font.name = "Consolas"
p.font.size = Pt(72)
p.font.bold = True
p.font.color.rgb = BORDER_COLOR
p.alignment = PP_ALIGN.CENTER

# Separator bar
p2 = tf.add_paragraph()
p2.text = "=================================================="
p2.font.name = "Consolas"
p2.font.size = Pt(18)
p2.font.color.rgb = GREEN_COLOR
p2.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide1.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(11.3), Inches(1.5))
tf_sub = subtitle_box.text_frame
tf_sub.word_wrap = True
p_sub = tf_sub.paragraphs[0]
p_sub.text = "Cyberpunk Student Programming Portal & Learning Console"
p_sub.font.name = "Consolas"
p_sub.font.size = Pt(20)
p_sub.font.bold = True
p_sub.font.color.rgb = WHITE_COLOR
p_sub.alignment = PP_ALIGN.CENTER

p_sub2 = tf_sub.add_paragraph()
p_sub2.text = "Simplified Outlines • Practical MCQ Quizzes • 3D Memory Cards • Free AI Tutor"
p_sub2.font.name = "Consolas"
p_sub2.font.size = Pt(14)
p_sub2.font.color.rgb = GRAY_COLOR
p_sub2.alignment = PP_ALIGN.CENTER


# ----------------- SLIDE 2: Problem Statement -----------------
slide2 = prs.slides.add_slide(slide_layout)
set_slide_background(slide2)
add_slide_title = add_slide_header(slide2, "> The Problem")

# 3 vertical cards side-by-side
card_width = Inches(3.64)
card_height = Inches(4.8)
top_pos = Inches(1.6)

draw_card(
    slide2, Inches(0.8), top_pos, card_width, card_height, 
    "01", "Boring Materials", 
    "Standard class notes and online books are static and dry. Students struggle to stay focused and rarely read them before coding tests."
)
draw_card(
    slide2, Inches(4.84), top_pos, card_width, card_height, 
    "02", "API Key Barriers", 
    "Most online AI helpers are complicated, require premium API keys, or get blocked by security checks in frontend browsers."
)
draw_card(
    slide2, Inches(8.88), top_pos, card_width, card_height, 
    "03", "No Study Logs", 
    "Standard class websites do not track learning progress, streaks, or student achievements, leading to low study consistency."
)


# ----------------- SLIDE 3: The CodeVerse Solution -----------------
slide3 = prs.slides.add_slide(slide_layout)
set_slide_background(slide3)
add_slide_header(slide3, "> The CodeVerse Solution")

# 2x2 grid cards matching user screenshot style
width_2x2 = Inches(5.7)
height_2x2 = Inches(2.3)

draw_card(
    slide3, Inches(0.8), Inches(1.6), width_2x2, height_2x2, 
    "01", "Study Tracks", 
    "Interactive topics sidebar showing page summaries and integrated note PDF handbooks in a clean portal."
)
draw_card(
    slide3, Inches(6.8), Inches(1.6), width_2x2, height_2x2, 
    "02", "Active Recall", 
    "Tactile 3D card deck flips and output prediction quizzes to build muscle memory and test syntax patterns."
)
draw_card(
    slide3, Inches(0.8), Inches(4.4), width_2x2, height_2x2, 
    "03", "Live AI Solver", 
    "Free, keyless AI assistant running client-side inside the browser to resolve doubts in real-time."
)
draw_card(
    slide3, Inches(6.8), Inches(4.4), width_2x2, height_2x2, 
    "04", "Session Auth", 
    "Cyberpunk login and registration console that secures study progress and caches developer logs."
)


# ----------------- SLIDE 4: Tech Stack (New Slide) -----------------
slide4 = prs.slides.add_slide(slide_layout)
set_slide_background(slide4)
add_slide_header(slide4, "> Programming Tech Stack")

# 2x2 Grid card display for HTML, CSS, JS, Python
draw_card(
    slide4, Inches(0.8), Inches(1.6), width_2x2, height_2x2, 
    "HTML", "Structure & Layout", 
    "Creates the webpage layouts, login forms, terminal frames, and inline SVG elements (circular dials and progress columns)."
)
draw_card(
    slide4, Inches(6.8), Inches(1.6), width_2x2, height_2x2, 
    "CSS", "Styles & 3D Effects", 
    "Builds the dark theme glows, glowing border layouts, responsive grid margins, and 3D card flipping animations."
)
draw_card(
    slide4, Inches(0.8), Inches(4.4), width_2x2, height_2x2, 
    "JS", "Application Logic", 
    "Manages quiz timers, flashcard keyboard hotkeys, localStorage progress caching, and live AI tutor connections."
)
draw_card(
    slide4, Inches(6.8), Inches(4.4), width_2x2, height_2x2, 
    "PY", "Local Dev Server", 
    "Spawns the local web server to bypass CORS blocks, runs note extractors, and serves PDF handbooks natively."
)


# ----------------- SLIDE 5: Tracks & Note PDF Integration -----------------
slide5 = prs.slides.add_slide(slide_layout)
set_slide_background(slide5)
add_slide_header(slide5, "> Dynamic Tracks & PDF Handbooks")

# 2 column layout
col_width = Inches(5.7)
col_height = Inches(4.8)

draw_card(
    slide5, Inches(0.8), Inches(1.6), col_width, col_height, 
    "01", "Interactive Docs", 
    "Displays summarized outlines fetched from Wikipedia and recommended books from Open Library. Dynamic track colors switch automatically based on selection."
)
draw_card(
    slide5, Inches(6.8), Inches(1.6), col_width, col_height, 
    "02", "PDF Note Viewer", 
    "A dual-tab toggle renders the professor's PDF handbook directly inside the browser using an iframe. Features single-click download actions."
)


# ----------------- SLIDE 6: 3D Flashcards -----------------
slide6 = prs.slides.add_slide(slide_layout)
set_slide_background(slide6)
add_slide_header(slide6, "> 3D Active Recall Flashcards")

draw_card(
    slide6, Inches(0.8), Inches(1.6), col_width, col_height, 
    "01", "Interactive Flip", 
    "Cards flip in 3D using mouse clicks or keyboard controls (Spacebar to flip, Arrow keys to navigate). Shuffling randomizes card orders."
)
draw_card(
    slide6, Inches(6.8), Inches(1.6), col_width, col_height, 
    "02", "Concept Databases", 
    "Dynamic generators scale the index pool to exactly 100 unique flashcards per track (400 total) with study validation tips."
)


# ----------------- SLIDE 7: Live AI Tutor -----------------
slide7 = prs.slides.add_slide(slide_layout)
set_slide_background(slide7)
add_slide_header(slide7, "> Live Online AI Solver")

draw_card(
    slide7, Inches(0.8), Inches(1.6), col_width, col_height, 
    "01", "Keyless API Connection", 
    "Powered by Pollinations AI. Runs live text inference directly inside the browser without asking the student for API keys."
)
draw_card(
    slide7, Inches(6.8), Inches(1.6), col_width, col_height, 
    "02", "Offline Database", 
    "Loaded with 100 interview-asked questions. Falls back to simulated answers instantly if the user is offline."
)


# ----------------- SLIDE 8: Analytics & Streaks -----------------
slide8 = prs.slides.add_slide(slide_layout)
set_slide_background(slide8)
add_slide_header(slide8, "> Student Progress Analytics")

draw_card(
    slide8, Inches(0.8), Inches(1.6), col_width, col_height, 
    "01", "Dashboard Metrics", 
    "Streaks card, SVG progress rings, and average quiz charts update in real-time, saved securely inside browser localStorage."
)
draw_card(
    slide8, Inches(6.8), Inches(1.6), col_width, col_height, 
    "02", "Compiler Achievements", 
    "Achievement badges (e.g. Card Master, Systems Engineer) unlock automatically when milestones are reached, displaying terminal toasts."
)


# ----------------- SLIDE 9: Conclusion -----------------
slide9 = prs.slides.add_slide(slide_layout)
set_slide_background(slide9)
add_slide_header(slide9, "> Project Conclusion", color=GREEN_COLOR)

conclusion_box = slide9.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.5))
tf9 = conclusion_box.text_frame
tf9.word_wrap = True

p = tf9.paragraphs[0]
p.text = "CodeVerse is a highly interactive learning portal that helps students learn coding in a gamified environment."
p.font.name = "Consolas"
p.font.size = Pt(22)
p.font.color.rgb = WHITE_COLOR
p.space_after = Pt(24)

p2 = tf9.add_paragraph()
p2.text = "✓ 400 MCQ Quiz Questions\n✓ 400 Concept Flashcards\n✓ 400 Interview Presets\n✓ Client-Side State Caching\n✓ 100% Free Live Online AI Tutor"
p2.font.name = "Consolas"
p2.font.size = Pt(20)
p2.font.bold = True
p2.font.color.rgb = BORDER_COLOR
p2.space_after = Pt(24)

p3 = tf9.add_paragraph()
p3.text = "[SYSTEM STATUS]: Presentation Compiled. Local Server Running."
p3.font.name = "Consolas"
p3.font.size = Pt(16)
p3.font.color.rgb = GREEN_COLOR

# Save the Presentation
output_path = os.path.join(r"C:\Users\SUCHISMITA\.gemini\antigravity\scratch\codeverse", "CodeVerse_Presentation_v2.pptx")
prs.save(output_path)
print("Updated presentation saved successfully at:", output_path)
