import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# --- Color Scheme: Realistic Green & Black Cyberpunk Terminal Theme ---
BG_COLOR = RGBColor(6, 6, 8)            # Deep Matte Black
CARD_BG_COLOR = RGBColor(16, 18, 16)     # Dark Terminal Gray-Black
BORDER_COLOR = RGBColor(0, 200, 80)      # Medium Terminal Green
GREEN_PRIMARY = RGBColor(0, 255, 102)    # Neon Cyber Green
GREEN_SECONDARY = RGBColor(140, 255, 140) # Matrix Soft Green
ACCENT_CYAN = RGBColor(0, 229, 255)      # Terminal Teal Accent
WHITE_TEXT = RGBColor(245, 245, 245)     # Primary White Text
GRAY_TEXT = RGBColor(160, 165, 160)      # Secondary Gray Text

# Helper function to set slide background
def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

# Helper function to draw header with a glowing line
def add_slide_header(slide, title_text, color=GREEN_PRIMARY):
    # Title Text Box
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.73), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Consolas"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = color
    
    # Glowing line under title (thin filled rectangle)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0.8), Inches(1.1), Inches(11.73), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER_COLOR
    line.line.fill.background() # No border line

# Helper function to draw a card panel with a divider and body text
def draw_card(slide, left, top, width, height, num_str, title_str, body_str, tag_str=None):
    # Base Card Shape
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG_COLOR
    card.line.color.rgb = BORDER_COLOR
    card.line.width = Pt(1.5)
    
    # Number & Title Text Box (Top of the card)
    title_box = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.2), width - Inches(0.6), Inches(0.5))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    
    # Number
    run_num = p_title.add_run()
    run_num.text = f"{num_str}   "
    run_num.font.name = "Consolas"
    run_num.font.bold = True
    run_num.font.size = Pt(18)
    run_num.font.color.rgb = GREEN_PRIMARY
    
    # Title
    run_title = p_title.add_run()
    run_title.text = title_str
    run_title.font.name = "Consolas"
    run_title.font.bold = True
    run_title.font.size = Pt(18)
    run_title.font.color.rgb = WHITE_TEXT
    
    # Optional Tag (Top Right)
    if tag_str:
        tag_box = slide.shapes.add_textbox(left + width - Inches(2.3), top + Inches(0.2), Inches(2.0), Inches(0.5))
        tf_tag = tag_box.text_frame
        tf_tag.word_wrap = True
        tf_tag.margin_left = tf_tag.margin_right = tf_tag.margin_top = tf_tag.margin_bottom = 0
        p_tag = tf_tag.paragraphs[0]
        p_tag.alignment = PP_ALIGN.RIGHT
        run_tag = p_tag.add_run()
        run_tag.text = tag_str
        run_tag.font.name = "Consolas"
        run_tag.font.bold = True
        run_tag.font.size = Pt(13)
        run_tag.font.color.rgb = ACCENT_CYAN
        
    # Divider Line inside Card
    div_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + Inches(0.3), top + Inches(0.7), width - Inches(0.6), Inches(0.015)
    )
    div_line.fill.solid()
    div_line.fill.fore_color.rgb = RGBColor(40, 60, 40)
    div_line.line.fill.background()
    
    # Body Text Box (below the divider)
    body_box = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.85), width - Inches(0.6), height - Inches(1.05))
    tf_body = body_box.text_frame
    tf_body.word_wrap = True
    tf_body.margin_left = tf_body.margin_right = tf_body.margin_top = tf_body.margin_bottom = 0
    
    p_body = tf_body.paragraphs[0]
    run_body = p_body.add_run()
    run_body.text = body_str
    run_body.font.name = "Consolas"
    run_body.font.size = Pt(13)
    run_body.font.color.rgb = GRAY_TEXT

# Helper to draw interactive keycaps
def draw_keycap(slide, left, top, width, height, key_name, action_desc):
    key = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    key.fill.solid()
    key.fill.fore_color.rgb = RGBColor(30, 36, 30)
    key.line.color.rgb = GREEN_PRIMARY
    key.line.width = Pt(1.5)
    
    tf = key.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = key_name
    p.font.name = "Consolas"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE_TEXT
    p.alignment = PP_ALIGN.CENTER
    
    desc_box = slide.shapes.add_textbox(left + width + Inches(0.25), top, Inches(4.8), height)
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    tf_desc.margin_left = tf_desc.margin_right = tf_desc.margin_top = tf_desc.margin_bottom = 0
    p_desc = tf_desc.paragraphs[0]
    p_desc.text = action_desc
    p_desc.font.name = "Consolas"
    p_desc.font.size = Pt(13)
    p_desc.font.color.rgb = GRAY_TEXT

# Helper to draw a mock chat dialogue bubble
def draw_chat_bubble(slide, left, top, width, height, sender, message, bg_color, border_color):
    bubble = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    bubble.fill.solid()
    bubble.fill.fore_color.rgb = bg_color
    bubble.line.color.rgb = border_color
    bubble.line.width = Pt(1.5)
    
    tf = bubble.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    
    p = tf.paragraphs[0]
    run_sender = p.add_run()
    run_sender.text = f"{sender}\n"
    run_sender.font.name = "Consolas"
    run_sender.font.bold = True
    run_sender.font.size = Pt(12)
    run_sender.font.color.rgb = border_color
    
    run_msg = p.add_run()
    run_msg.text = message
    run_msg.font.name = "Consolas"
    run_msg.font.size = Pt(11)
    run_msg.font.color.rgb = WHITE_TEXT

# ----------------- SLIDE 1: Title Slide (Boot Terminal Theme) -----------------
slide_layout = prs.slide_layouts[6] # Blank
slide1 = prs.slides.add_slide(slide_layout)
set_slide_background(slide1)

title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(1.8))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "CodeVerse"
p.font.name = "Consolas"
p.font.size = Pt(72)
p.font.bold = True
p.font.color.rgb = GREEN_PRIMARY
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "=================================================="
p2.font.name = "Consolas"
p2.font.size = Pt(18)
p2.font.color.rgb = BORDER_COLOR
p2.alignment = PP_ALIGN.CENTER

subtitle_box = slide1.shapes.add_textbox(Inches(1.0), Inches(3.8), Inches(11.3), Inches(1.2))
tf_sub = subtitle_box.text_frame
tf_sub.word_wrap = True
p_sub = tf_sub.paragraphs[0]
p_sub.text = "Cyberpunk Active-Recall Study Console & Interactive Portal"
p_sub.font.name = "Consolas"
p_sub.font.size = Pt(20)
p_sub.font.bold = True
p_sub.font.color.rgb = WHITE_TEXT
p_sub.alignment = PP_ALIGN.CENTER

# Boot sequence logs (Realistic human touch)
logs_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(11.3), Inches(1.8))
tf_logs = logs_box.text_frame
tf_logs.word_wrap = True
p_log = tf_logs.paragraphs[0]
p_log.alignment = PP_ALIGN.CENTER
p_log.text = "[OK] Loaded: 4 language tracks | [OK] Configured: 400 interactive flashcards\n[OK] Mounted: Free Pollinations AI Doubt Engine (Keyless Connection)\n[OK] Local Host: system_init_session.sh -> PORT 8081 ready"
p_log.font.name = "Consolas"
p_log.font.size = Pt(12)
p_log.font.color.rgb = GRAY_TEXT


# ----------------- SLIDE 2: The Problem (Asymmetric Layout) -----------------
slide2 = prs.slides.add_slide(slide_layout)
set_slide_background(slide2)
add_slide_header(slide2, "> The Problem")

# Left Highlight Card (Large)
draw_card(
    slide2, Inches(0.8), Inches(1.6), Inches(5.4), Inches(4.8), 
    "01", "Boring Materials", 
    "Traditional textbooks and lecture notes are static, dry, and boring. Students struggle to read them and lose interest quickly before programming tests."
)

# Right Stacked Cards (Two smaller horizontal cards)
draw_card(
    slide2, Inches(6.8), Inches(1.6), Inches(5.73), Inches(2.25), 
    "02", "API Key Barriers", 
    "Most AI coding assistants demand premium API keys or complex environment variables, blocking students from instant help."
)
draw_card(
    slide2, Inches(6.8), Inches(4.15), Inches(5.73), Inches(2.25), 
    "03", "No Study Logs", 
    "Standard school websites do not track learning habits, active daily streaks, or practice achievements."
)


# ----------------- SLIDE 3: The CodeVerse Solution (3 columns) -----------------
slide3 = prs.slides.add_slide(slide_layout)
set_slide_background(slide3)
add_slide_header(slide3, "> The CodeVerse Solution")

col_w = Inches(3.64)
col_h = Inches(4.8)
top_pos = Inches(1.6)

draw_card(
    slide3, Inches(0.8), top_pos, col_w, col_h, 
    "01", "Interactive Docs", 
    "Integrates professor PDF notes side-by-side with Wikipedia outlines and Open Library links inside a single responsive browser console."
)
draw_card(
    slide3, Inches(4.84), top_pos, col_w, col_h, 
    "02", "Active Recall", 
    "Replaces passive reading with 3D concept flashcards and output prediction quizzes that shuffle questions dynamically."
)
draw_card(
    slide3, Inches(8.88), top_pos, col_w, col_h, 
    "03", "Keyless AI Solver", 
    "Connects to Pollinations AI to resolve doubts instantly for free directly in the browser—no accounts or API keys required."
)


# ----------------- SLIDE 4: Programming Tech Stack (Tag-based Grid) -----------------
slide4 = prs.slides.add_slide(slide_layout)
set_slide_background(slide4)
add_slide_header(slide4, "> Programming Tech Stack")

width_2x2 = Inches(5.7)
height_2x2 = Inches(2.3)

draw_card(
    slide4, Inches(0.8), Inches(1.6), width_2x2, height_2x2, 
    "HTML", "Structure & Layout", 
    "Builds all 7 portal pages, establishes the terminal boxes, interactive forms, and embeds inline SVGs for circular progress dials.",
    tag_str="[MARKUP]"
)
draw_card(
    slide4, Inches(6.8), Inches(1.6), width_2x2, height_2x2, 
    "CSS", "Styling & 3D Effects", 
    "Controls the Obsidian-green dark theme, glow effects, responsive grid alignments, and fluid 3D card-flip animations.",
    tag_str="[DESIGN]"
)
draw_card(
    slide4, Inches(0.8), Inches(4.4), width_2x2, height_2x2, 
    "JS", "Application Logic", 
    "Handles session streaks, quiz timers, keyboard hooks, localStorage state caching, and the keyless AI tutor API stream.",
    tag_str="[CLIENT]"
)
draw_card(
    slide4, Inches(6.8), Inches(4.4), width_2x2, height_2x2, 
    "PY", "Dev Server & Scripts", 
    "Launches the local HTTP server to host PDF handbooks natively in the browser and runs scripts to verify note file structures.",
    tag_str="[UTILITY]"
)


# ----------------- SLIDE 5: How CodeVerse Works (Interactive Pipeline) -----------------
slide5 = prs.slides.add_slide(slide_layout)
set_slide_background(slide5)
add_slide_header(slide5, "> How CodeVerse Works")

step_w = Inches(3.4)
step_h = Inches(4.5)
step_top = Inches(1.8)

# Step 1
draw_card(
    slide5, Inches(0.8), step_top, step_w, step_h, 
    "STEP 01", "Sign In", 
    "Students log in via the terminal UI (or bypass as Guest) to initialize their workspace and restore study streaks."
)

# Arrow Connector 1
arrow1_box = slide5.shapes.add_textbox(Inches(4.3), Inches(3.6), Inches(0.5), Inches(0.5))
p_a1 = arrow1_box.text_frame.paragraphs[0]
p_a1.text = ">>>"
p_a1.font.name = "Consolas"
p_a1.font.bold = True
p_a1.font.size = Pt(20)
p_a1.font.color.rgb = GREEN_PRIMARY

# Step 2
draw_card(
    slide5, Inches(4.9), step_top, step_w, step_h, 
    "STEP 02", "Study & Practice", 
    "Pick a track, review outlines side-by-side with PDF handbooks, flip flashcards, and run MCQ quiz engines to practice."
)

# Arrow Connector 2
arrow2_box = slide5.shapes.add_textbox(Inches(8.4), Inches(3.6), Inches(0.5), Inches(0.5))
p_a2 = arrow2_box.text_frame.paragraphs[0]
p_a2.text = ">>>"
p_a2.font.name = "Consolas"
p_a2.font.bold = True
p_a2.font.size = Pt(20)
p_a2.font.color.rgb = GREEN_PRIMARY

# Step 3
draw_card(
    slide5, Inches(9.0), step_top, step_w, step_h, 
    "STEP 03", "Resolve & Track", 
    "Ask the AI tutor about syntax errors. Progress metrics, streaks, and achievements are auto-saved to localStorage."
)


# ----------------- SLIDE 6: Tracks & Note PDF Integration (Split Layout) -----------------
slide6 = prs.slides.add_slide(slide_layout)
set_slide_background(slide6)
add_slide_header(slide6, "> Dynamic Tracks & PDF Handbooks")

col_width = Inches(5.7)
col_height = Inches(4.8)

draw_card(
    slide6, Inches(0.8), Inches(1.6), col_width, col_height, 
    "01", "Interactive Outlines", 
    "Browse topics, Wikipedia details, and books. Page accents dynamically swap accent colors (Teal, Amber, Violet, Green) to match the selected language track."
)
draw_card(
    slide6, Inches(6.8), Inches(1.6), col_width, col_height, 
    "02", "Embedded PDF Viewer", 
    "Toggle from summaries to original professor PDF handbooks directly inside the browser using custom terminal iframe containers."
)


# ----------------- SLIDE 7: 3D Flashcards (Visual Keyboard Vibe) -----------------
slide7 = prs.slides.add_slide(slide_layout)
set_slide_background(slide7)
add_slide_header(slide7, "> 3D Active Recall Flashcards")

draw_card(
    slide7, Inches(0.8), Inches(1.6), col_width, col_height, 
    "01", "Holographic Deck", 
    "Cards flip 180° in 3D using mouse clicks or hotkeys. Shuffle deck to randomize order and tag cards as 'Mastered' to log study data."
)

# Right side: Simulated Keyboard keys (Realistic Human Touch)
guide_box = slide7.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(5.7), Inches(0.5))
p_g = guide_box.text_frame.paragraphs[0]
p_g.text = "[ ACTIVE KEYBOARD SHORTCUTS ]"
p_g.font.name = "Consolas"
p_g.font.bold = True
p_g.font.size = Pt(16)
p_g.font.color.rgb = GREEN_PRIMARY

draw_keycap(slide7, Inches(6.8), Inches(2.4), Inches(1.8), Inches(0.6), "SPACEBAR", "Flips the active card to view answer/compiler tips")
draw_keycap(slide7, Inches(6.8), Inches(3.4), Inches(1.8), Inches(0.6), "ARROW LEFT", "Navigates back to the previous flashcard")
draw_keycap(slide7, Inches(6.8), Inches(4.4), Inches(1.8), Inches(0.6), "ARROW RIGHT", "Navigates forward to the next flashcard")


# ----------------- SLIDE 8: Live AI Tutor (Simulated Chat UI) -----------------
slide8 = prs.slides.add_slide(slide_layout)
set_slide_background(slide8)
add_slide_header(slide8, "> Live Online AI Solver")

draw_card(
    slide8, Inches(0.8), Inches(1.6), col_width, col_height, 
    "01", "Pollinations AI", 
    "Queries are sent directly to the AI model on the client side. Operates completely keyless and falls back to 100 pre-saved interview questions if offline."
)

# Right side: Mock Chat Interface
draw_chat_bubble(
    slide8, Inches(6.8), Inches(1.8), Inches(5.7), Inches(1.5), 
    "student@codeverse:~$ Explain CSS Grid auto-fit", 
    "How does the auto-fit keyword work in CSS Grid?",
    RGBColor(24, 30, 24), GREEN_PRIMARY
)

draw_chat_bubble(
    slide8, Inches(6.8), Inches(3.6), Inches(5.7), Inches(2.4), 
    "[AI_TUTOR_BOT] @ Pollinations_Inference", 
    "The auto-fit keyword stretches existing grid items to fill all available space in the track, preventing empty gaps on wider viewports.\n\nExample:\ngrid-template-columns: repeat(auto-fit, minmax(200px, 1fr));",
    RGBColor(16, 18, 16), ACCENT_CYAN
)


# ----------------- SLIDE 9: Analytics & Streaks (Simulated Dashboard) -----------------
slide9 = prs.slides.add_slide(slide_layout)
set_slide_background(slide9)
add_slide_header(slide9, "> Student Progress Analytics")

draw_card(
    slide9, Inches(0.8), Inches(1.6), col_width, col_height, 
    "01", "Dashboard Metrics", 
    "Displays study streaks, SVG completion dials, and bar charts of quiz scores saved locally. Unlocks badges for reaching key study milestones."
)

# Right side: Simulated dashboard badges & progress bars
dash_card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.7), Inches(4.8))
dash_card.fill.solid()
dash_card.fill.fore_color.rgb = CARD_BG_COLOR
dash_card.line.color.rgb = BORDER_COLOR
dash_card.line.width = Pt(1.5)

dash_box = slide9.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.1), Inches(4.4))
tf_dash = dash_box.text_frame
tf_dash.word_wrap = True
tf_dash.margin_left = tf_dash.margin_right = tf_dash.margin_top = tf_dash.margin_bottom = 0

p_dash = tf_dash.paragraphs[0]
p_dash.text = "[ ACHIEVEMENTS LOADED ]"
p_dash.font.name = "Consolas"
p_dash.font.bold = True
p_dash.font.size = Pt(16)
p_dash.font.color.rgb = GREEN_PRIMARY
p_dash.space_after = Pt(14)

p_badges = tf_dash.add_paragraph()
p_badges.text = "★ [SYSTEM ENGINEER] - Logged a 5-day study streak\n★ [CARD SCHOLAR]   - Mastered 10+ flashcards\n★ [COMPILER MASTER] - Scored 100% on MCQ Quiz"
p_badges.font.name = "Consolas"
p_badges.font.bold = True
p_badges.font.size = Pt(12)
p_badges.font.color.rgb = WHITE_TEXT
p_badges.space_after = Pt(24)

p_tracks_title = tf_dash.add_paragraph()
p_tracks_title.text = "[ TRACK COMPLETION METRICS ]"
p_tracks_title.font.name = "Consolas"
p_tracks_title.font.bold = True
p_tracks_title.font.size = Pt(14)
p_tracks_title.font.color.rgb = ACCENT_CYAN
p_tracks_title.space_after = Pt(10)

p_bars = tf_dash.add_paragraph()
p_bars.text = "Python Track: [====================] 100%\nJS Track:     [==============      ] 70%\nC++ Track:    [==========          ] 50%\nDSA Track:    [======              ] 30%"
p_bars.font.name = "Consolas"
p_bars.font.size = Pt(12)
p_bars.font.color.rgb = GRAY_TEXT


# ----------------- SLIDE 10: Conclusion (System Log Output) -----------------
slide10 = prs.slides.add_slide(slide_layout)
set_slide_background(slide10)
add_slide_header(slide10, "> Project Conclusion")

conclusion_box = slide10.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.33), Inches(4.5))
tf10 = conclusion_box.text_frame
tf10.word_wrap = True
tf10.margin_left = tf10.margin_right = tf10.margin_top = tf10.margin_bottom = 0

p = tf10.paragraphs[0]
p.text = "CodeVerse successfully bridges the gap between offline notes and active programming practice. It delivers a modern, gamified study console to make coding accessible."
p.font.name = "Consolas"
p.font.size = Pt(20)
p.font.color.rgb = WHITE_TEXT
p.space_after = Pt(20)

p2 = tf10.add_paragraph()
p2.text = "[CORE CONSOLE METRICS STATUS]\n-------------------------------------------------------------\n✓ 400 MCQ Quiz Questions  | ✓ 400 Active Recall Flashcards\n✓ 400 Interview Presets    | ✓ 100% Free Live Online AI Tutor\n✓ Offline Storage Cache   | ✓ Responsive Terminal Framework"
p2.font.name = "Consolas"
p2.font.bold = True
p2.font.size = Pt(16)
p2.font.color.rgb = GREEN_SECONDARY
p2.space_after = Pt(24)

p3 = tf10.add_paragraph()
p3.text = "[SYSTEM LOG]: Deployment Complete. Server running successfully on port 8081."
p3.font.name = "Consolas"
p3.font.size = Pt(16)
p3.font.color.rgb = GREEN_PRIMARY

# Save the Presentation
output_path = os.path.join(r"C:\Users\SUCHISMITA\.gemini\antigravity\scratch\codeverse", "CodeVerse_Presentation_v3.pptx")
prs.save(output_path)
print("Updated presentation saved successfully at:", output_path)
