import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# --- Color Scheme: Purple Dark & Green Dark Theme ---
BG_COLOR = RGBColor(10, 8, 20)          # Very Dark Obsidian Purple
CARD_BG_COLOR = RGBColor(22, 18, 38)     # Slightly Lighter Dark Purple
BORDER_COLOR = RGBColor(155, 93, 229)    # Neon Purple/Violet (Primary Accent)
GREEN_COLOR = RGBColor(0, 245, 150)      # Cyber Green (Secondary Accent)
WHITE_COLOR = RGBColor(255, 255, 255)    # Text Primary
GRAY_COLOR = RGBColor(175, 172, 190)     # Text Secondary

# Helper function to set slide background
def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

# Helper function to draw header with a glowing line
def add_slide_header(slide, title_text, color=GREEN_COLOR):
    # Title Text Box
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.73), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Consolas"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = color
    
    # Glowing line under title (thin filled rectangle)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0.8), Inches(1.1), Inches(11.73), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER_COLOR
    line.line.fill.background() # No border line

# Helper function to draw a card panel with a divider and body text
def draw_card(slide, left, top, width, height, num_str, title_str, body_str):
    # Base Card Shape
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG_COLOR
    card.line.color.rgb = BORDER_COLOR
    card.line.width = Pt(1.5)
    
    # Number & Title Text Box (Top of the card)
    title_box = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.2), width - Inches(0.6), Inches(0.5))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    
    # Number
    run_num = p_title.add_run()
    run_num.text = f"{num_str}   "
    run_num.font.name = "Consolas"
    run_num.font.bold = True
    run_num.font.size = Pt(18)
    run_num.font.color.rgb = GREEN_COLOR
    
    # Title
    run_title = p_title.add_run()
    run_title.text = title_str
    run_title.font.name = "Consolas"
    run_title.font.bold = True
    run_title.font.size = Pt(18)
    run_title.font.color.rgb = WHITE_COLOR
    
    # Divider Line inside Card (neat thin rectangle instead of dashes)
    div_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + Inches(0.3), top + Inches(0.7), width - Inches(0.6), Inches(0.015)
    )
    div_line.fill.solid()
    div_line.fill.fore_color.rgb = RGBColor(60, 50, 80)
    div_line.line.fill.background()
    
    # Body Text Box (below the divider)
    body_box = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.85), width - Inches(0.6), height - Inches(1.05))
    tf_body = body_box.text_frame
    tf_body.word_wrap = True
    tf_body.margin_left = tf_body.margin_right = tf_body.margin_top = tf_body.margin_bottom = 0
    
    p_body = tf_body.paragraphs[0]
    run_body = p_body.add_run()
    run_body.text = body_str
    run_body.font.name = "Consolas"
    run_body.font.size = Pt(13)
    run_body.font.color.rgb = GRAY_COLOR

# ----------------- SLIDE 1: Title Slide -----------------
slide_layout = prs.slide_layouts[6] # Blank
slide1 = prs.slides.add_slide(slide_layout)
set_slide_background(slide1)

# Large Title
title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.0))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "CodeVerse"
p.font.name = "Consolas"
p.font.size = Pt(72)
p.font.bold = True
p.font.color.rgb = BORDER_COLOR
p.alignment = PP_ALIGN.CENTER

# Separator bar
p2 = tf.add_paragraph()
p2.text = "=================================================="
p2.font.name = "Consolas"
p2.font.size = Pt(18)
p2.font.color.rgb = GREEN_COLOR
p2.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide1.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(11.3), Inches(1.5))
tf_sub = subtitle_box.text_frame
tf_sub.word_wrap = True
p_sub = tf_sub.paragraphs[0]
p_sub.text = "Cyberpunk Student Programming Portal & Learning Console"
p_sub.font.name = "Consolas"
p_sub.font.size = Pt(20)
p_sub.font.bold = True
p_sub.font.color.rgb = WHITE_COLOR
p_sub.alignment = PP_ALIGN.CENTER

p_sub2 = tf_sub.add_paragraph()
p_sub2.text = "Simplified Outlines • Practical MCQ Quizzes • 3D Memory Cards • Free AI Tutor"
p_sub2.font.name = "Consolas"
p_sub2.font.size = Pt(14)
p_sub2.font.color.rgb = GRAY_COLOR
p_sub2.alignment = PP_ALIGN.CENTER


# ----------------- SLIDE 2: Problem Statement -----------------
slide2 = prs.slides.add_slide(slide_layout)
set_slide_background(slide2)
add_slide_header(slide2, "> The Problem")

# 3 vertical cards side-by-side
card_width = Inches(3.64)
card_height = Inches(4.8)
top_pos = Inches(1.6)

draw_card(
    slide2, Inches(0.8), top_pos, card_width, card_height, 
    "01", "Boring Notes", 
    "Class notes and coding textbooks are dry and boring. Students rarely read them before coding tests."
)
draw_card(
    slide2, Inches(4.84), top_pos, card_width, card_height, 
    "02", "AI Key Issues", 
    "Most AI coding assistants are complex to set up or require paid API keys, making them hard for students to use."
)
draw_card(
    slide2, Inches(8.88), top_pos, card_width, card_height, 
    "03", "No Study Logs", 
    "Standard study portals do not track daily habits, practice consistency, or student achievements."
)


# ----------------- SLIDE 3: The CodeVerse Solution -----------------
slide3 = prs.slides.add_slide(slide_layout)
set_slide_background(slide3)
add_slide_header(slide3, "> The CodeVerse Solution")

# 2x2 grid cards matching user screenshot style
width_2x2 = Inches(5.7)
height_2x2 = Inches(2.3)

draw_card(
    slide3, Inches(0.8), Inches(1.6), width_2x2, height_2x2, 
    "01", "Study Tracks", 
    "Structured learning modules with summaries and built-in PDF handbooks in a clean console."
)
draw_card(
    slide3, Inches(6.8), Inches(1.6), width_2x2, height_2x2, 
    "02", "Active Recall", 
    "Interactive 3D flashcards and output prediction quizzes to build muscle memory."
)
draw_card(
    slide3, Inches(0.8), Inches(4.4), width_2x2, height_2x2, 
    "03", "Live AI Solver", 
    "Free, keyless AI chatbot running directly in the browser to answer coding questions instantly."
)
draw_card(
    slide3, Inches(6.8), Inches(4.4), width_2x2, height_2x2, 
    "04", "Session Auth", 
    "A secure login system that saves your personal study progress, badges, and streaks."
)


# ----------------- SLIDE 4: Tech Stack (New Slide) -----------------
slide4 = prs.slides.add_slide(slide_layout)
set_slide_background(slide4)
add_slide_header(slide4, "> Programming Tech Stack")

# 2x2 Grid card display for HTML, CSS, JS, Python
draw_card(
    slide4, Inches(0.8), Inches(1.6), width_2x2, height_2x2, 
    "HTML", "Structure & Layout", 
    "Creates the framework for all 7 portal pages. Embeds terminal panels, forms, and SVGs for visual progress metrics."
)
draw_card(
    slide4, Inches(6.8), Inches(1.6), width_2x2, height_2x2, 
    "CSS", "Styles & 3D Effects", 
    "Builds the dark theme, neon glows, responsive layouts, and smooth 3D card-flip animations."
)
draw_card(
    slide4, Inches(0.8), Inches(4.4), width_2x2, height_2x2, 
    "JS", "Application Logic", 
    "Handles session streaks, quiz timers, keyboard hotkeys, state persistence, and live AI tutor connections."
)
draw_card(
    slide4, Inches(6.8), Inches(4.4), width_2x2, height_2x2, 
    "PY", "Dev Server & Scripts", 
    "Runs the local development server to load PDF handbooks natively and scan note file structures."
)


# ----------------- SLIDE 5: Tracks & Note PDF Integration -----------------
slide5 = prs.slides.add_slide(slide_layout)
set_slide_background(slide5)
add_slide_header(slide5, "> Dynamic Tracks & PDF Handbooks")

# 2 column layout
col_width = Inches(5.7)
col_height = Inches(4.8)

draw_card(
    slide5, Inches(0.8), Inches(1.6), col_width, col_height, 
    "01", "Interactive Outlines", 
    "Browse coding topics, Wikipedia summaries, and recommended books. Page accents dynamically shift colors to match the language selected."
)
draw_card(
    slide5, Inches(6.8), Inches(1.6), col_width, col_height, 
    "02", "Embedded PDF Viewer", 
    "Toggle between outlines and full textbook PDF handbooks. Reads documents directly in the browser using custom iframe windows."
)


# ----------------- SLIDE 6: 3D Flashcards -----------------
slide6 = prs.slides.add_slide(slide_layout)
set_slide_background(slide6)
add_slide_header(slide6, "> 3D Active Recall Flashcards")

draw_card(
    slide6, Inches(0.8), Inches(1.6), col_width, col_height, 
    "01", "Holographic Flips", 
    "Cards flip 180° in 3D. Control them using your keyboard (Space to flip, Left/Right arrows to navigate). Shuffling randomizes card order."
)
draw_card(
    slide6, Inches(6.8), Inches(1.6), col_width, col_height, 
    "02", "Concept Database", 
    "Features 100 flashcards per language track (400 total) to test syntax, with instant review tips on the card backs."
)


# ----------------- SLIDE 7: Live AI Tutor -----------------
slide7 = prs.slides.add_slide(slide_layout)
set_slide_background(slide7)
add_slide_header(slide7, "> Live Online AI Solver")

draw_card(
    slide7, Inches(0.8), Inches(1.6), col_width, col_height, 
    "01", "Keyless AI Tutor", 
    "Powered by Pollinations AI. Students can ask coding questions and get help without signing up or using API keys."
)
draw_card(
    slide7, Inches(6.8), Inches(1.6), col_width, col_height, 
    "02", "Offline Database", 
    "Includes 100 pre-saved interview questions. Falls back to simulated answers instantly if the user is offline."
)


# ----------------- SLIDE 8: Analytics & Streaks -----------------
slide8 = prs.slides.add_slide(slide_layout)
set_slide_background(slide8)
add_slide_header(slide8, "> Student Progress Analytics")

draw_card(
    slide8, Inches(0.8), Inches(1.6), col_width, col_height, 
    "01", "Metrics Dashboard", 
    "Displays daily active study streaks, circular SVG progress dials, and bar charts of your quiz scores. Saved locally in your browser."
)
draw_card(
    slide8, Inches(6.8), Inches(1.6), col_width, col_height, 
    "02", "Milestone Achievements", 
    "Unlock developer badges like 'Systems Engineer' or 'Card Master' for reaching study milestones, with toast notifications."
)


# ----------------- SLIDE 9: Conclusion -----------------
slide9 = prs.slides.add_slide(slide_layout)
set_slide_background(slide9)
add_slide_header(slide9, "> Project Conclusion", color=GREEN_COLOR)

conclusion_box = slide9.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.33), Inches(4.5))
tf9 = conclusion_box.text_frame
tf9.word_wrap = True
tf9.margin_left = tf9.margin_right = tf9.margin_top = tf9.margin_bottom = 0

p = tf9.paragraphs[0]
p.text = "CodeVerse bridges the gap between offline notes and active coding practice. It is a 100% client-side, gamified portal designed to help students learn coding easily."
p.font.name = "Consolas"
p.font.size = Pt(22)
p.font.color.rgb = WHITE_COLOR
p.space_after = Pt(24)

p2 = tf9.add_paragraph()
p2.text = "✓ 400 MCQ Quiz Questions\n✓ 400 Active Recall Flashcards\n✓ 100% Free AI Tutor (No API Key Required)\n✓ Interactive Daily Streak Tracker\n✓ Clean Cyberpunk Dark UI"
p2.font.name = "Consolas"
p2.font.size = Pt(20)
p2.font.bold = True
p2.font.color.rgb = BORDER_COLOR
p2.space_after = Pt(24)

p3 = tf9.add_paragraph()
p3.text = "[SYSTEM STATUS]: Presentation Compiled. Local Server Running."
p3.font.name = "Consolas"
p3.font.size = Pt(16)
p3.font.color.rgb = GREEN_COLOR

# Save the Presentation
output_path = os.path.join(r"C:\Users\SUCHISMITA\.gemini\antigravity\scratch\codeverse", "CodeVerse_Presentation_v2.pptx")
prs.save(output_path)
print("Updated presentation saved successfully at:", output_path)
